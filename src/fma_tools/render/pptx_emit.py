"""HTML deck to REAL PowerPoint: text boxes and native tables, never pictures.

The mechanism: Playwright lays the deck out, a JS walk reads every element's computed
position and style off the layout Chrome already solved once, and python-pptx emits
text boxes and native tables at those coordinates. Nothing here re-authors a layout,
and nothing estimates from character counts. Charts (svg, canvas, img, and anything
marked data-render="image") stay images -- an SVG re-drawn as shapes is a redrawing,
and a redrawing can differ from the figure the PDF carries.

Why native text is non-negotiable: a deck once shipped twice (25 then 30 slides) as
rasterised pictures of slides, to a recipient who finishes everything in PowerPoint
and had asked to add and delete pages. The delivered-bytes gates in gates.py make
that artifact structurally unshippable.

Measured constants, carried whole -- never re-derived, never widened:
  PPTX_ROW_GROWTH = 1.30  PowerPoint renders ~31px rows where Chrome renders 24px; a
                          row height is a MINIMUM everywhere except the browser. If a
                          table does not fit, the table changes, never this constant.
  LINE_SLACK = 1.18       shrink-to-fit single-line boxes widen for font substitution
                          (measured against LibreOffice, the honest worst case).
Line spacing is set in points, never a multiple (a 1.4 multiple renders nearer 1.7 in
PowerPoint); letter-spacing is a raw rPr attribute python-pptx has no property for.
"""

from __future__ import annotations

from pathlib import Path

from ..errors import Refusal

PPTX_ROW_GROWTH = 1.30
LINE_SLACK = 1.18

# PowerPoint needs the face installed on the OPENING machine, so only faces every
# Mac and Windows box carries are allowed; embedding is not supported by python-pptx.
SAFE_FONTS = ("Arial", "Calibri", "Georgia", "Times New Roman", "Verdana",
              "Courier New")

_MEASURE_JS = r"""() => {
  const slides = [...document.querySelectorAll('section.slide')];
  const problems = [];
  let imgIdx = 0;
  const out = [];
  const IMG = 'img, svg, canvas, [data-render="image"]';
  for (let si = 0; si < slides.length; si++) {
    const s = slides[si];
    const sr = s.getBoundingClientRect();
    const geo = (el) => { const r = el.getBoundingClientRect();
      return {x: r.left - sr.left, y: r.top - sr.top, w: r.width, h: r.height}; };
    const num = (v, fb) => { const n = parseFloat(v); return isNaN(n) ? fb : n; };
    const styleOf = (el) => { const cs = getComputedStyle(el);
      const size = parseFloat(cs.fontSize);
      return {size, lh: num(cs.lineHeight, size * 1.2),
              weight: parseInt(cs.fontWeight) || 400, colour: cs.color,
              align: cs.textAlign, transform: cs.textTransform,
              spacing: num(cs.letterSpacing, 0)}; };
    const texts = [], tables = [], images = [], fills = [], oob = [];

    const cellOf = (td) => { const cs = getComputedStyle(td); const g = geo(td);
      if (td.querySelector(IMG) || td.querySelector('table'))
        problems.push('slide ' + (si + 1) + ': a table cell contains a chart or ' +
          'nested table, which the PPTX table cannot carry; move it out of the ' +
          'table or restate it');
      return {x: g.x, w: g.w, text: td.textContent.trim(),
              align: cs.textAlign, weight: parseInt(cs.fontWeight) || 400,
              size: parseFloat(cs.fontSize), colour: cs.color,
              transform: cs.textTransform,
              fill: cs.backgroundColor}; };
    const tableOf = (t) => ({...geo(t),
      rows: [...t.querySelectorAll('tr')].filter(tr => tr.closest('table') === t)
        .map(tr => { const g = geo(tr);
        return {y: g.y, h: g.h, head: tr.querySelector('th') !== null,
                cells: [...tr.children].map(cellOf)}; })});

    const walk = (el) => {
      const cs = getComputedStyle(el);
      if (cs.display === 'none' || cs.visibility === 'hidden') return;
      if (el.matches(IMG)) {
        el.setAttribute('data-fma-img', String(imgIdx));
        images.push({...geo(el), idx: imgIdx++});
        return;
      }
      if (el.tagName.toLowerCase() === 'table') { tables.push(tableOf(el)); return; }
      const bg = cs.backgroundColor;
      if (bg && bg !== 'rgba(0, 0, 0, 0)' && bg !== 'transparent' && el !== s) {
        const g = geo(el);
        if (g.w > 0 && g.h > 0) fills.push({...g, colour: bg});
      }
      const own = [...el.childNodes].filter(n => n.nodeType === 3)
        .map(n => n.textContent).join('').trim();
      const kids = [...el.children];
      const kidHasText = kids.some(c => c.textContent.trim());
      // a text leaf must hold NO embedded chart or table anywhere beneath it, or
      // the graphic would be flattened into a run and vanish -- and text parity
      // would still pass, because 'expected' comes from this same walk
      const embedded = el.querySelector('table, ' + IMG) !== null;
      const allInline = kids.length > 0 && kids.every(c => {
        const d = getComputedStyle(c).display;
        return d.startsWith('inline') || c.tagName.toLowerCase() === 'br'; });
      const leaf = !embedded && ((own && !kidHasText)
        || (own && kidHasText && allInline)
        || (!own && kidHasText && allInline));
      if (leaf) {
        const g = geo(el);
        const pw = el.parentElement.getBoundingClientRect().width;
        const st = styleOf(el);
        texts.push({...g, ...st, text: el.textContent.trim(),
                    full: g.w >= pw - 1,
                    lines: Math.max(1, Math.round(g.h / st.lh))});
        return;
      }
      for (const c of kids) walk(c);
      // reaching here with own text means it was NOT captured as a leaf (block
      // children, or an embedded chart/table beneath) -- refuse, never drop
      if (own)
        problems.push('slide ' + (si + 1) + ': <' + el.tagName.toLowerCase() +
          '> mixes its own text with block children or an embedded chart; wrap ' +
          'the bare text in its own element so nothing is silently dropped: ' +
          JSON.stringify(own.slice(0, 60)));
    };
    const rootText = [...s.childNodes].filter(n => n.nodeType === 3)
      .map(n => n.textContent).join('').trim();
    if (rootText)
      problems.push('slide ' + (si + 1) + ': bare text sits directly inside ' +
        'section.slide and would be dropped -- wrap it in an element: ' +
        JSON.stringify(rootText.slice(0, 60)));
    for (const c of [...s.children]) walk(c);

    // bounds: no painted element may cross its slide box in any direction
    for (const el of s.querySelectorAll('*')) {
      const r = el.getBoundingClientRect();
      if (r.height <= 0 || r.width <= 0) continue;
      if (r.top < sr.top - 1 || r.left < sr.left - 1 ||
          r.bottom > sr.bottom + 1 || r.right > sr.right + 1)
        oob.push('slide ' + (si + 1) + ': <' + el.tagName.toLowerCase() + '> ' +
                 (el.className && typeof el.className === 'string'
                    ? '.' + el.className.split(' ')[0] : '') +
                 ' paints outside the slide box');
    }
    out.push({w: sr.width, h: sr.height, texts, tables, images, fills,
              out_of_bounds: [...new Set(oob)].slice(0, 6)});
  }
  return {slides: out, problems};
}"""


def _px(n: float):
    from pptx.util import Emu
    return Emu(int(round(n * 9525)))


def _pt(css_px: float):
    from pptx.util import Pt
    return Pt(round(css_px * 0.75, 1))


def _rgb(css: str):
    from pptx.dml.color import RGBColor
    css = css.strip()
    if css.startswith("#"):
        h = css[1:]
        if len(h) == 3:
            h = "".join(c * 2 for c in h)
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    if css.startswith(("rgba(", "rgb(")):
        parts = [p.strip() for p in css[css.index("(") + 1:css.rindex(")")].split(",")]
        r, g, b = (float(p) for p in parts[:3])
        a = float(parts[3]) if len(parts) > 3 else 1.0
        # composited against white: a translucent rule is a different grey in every
        # viewer, and hairlines are structure, not decoration
        return RGBColor(*(int(round(c * a + 255 * (1 - a))) for c in (r, g, b)))
    return RGBColor(0, 0, 0)


def shown(t: dict) -> str:
    """What the browser PAINTS: text-transform applied, so the .pptx and the page
    cannot drift apart on case."""
    return t["text"].upper() if t.get("transform") == "uppercase" else t["text"]


def measure(session, font: str) -> list[dict]:
    """Read every box off the loaded deck, under the face PowerPoint will use."""
    session.override_font(font)
    m = session.page.evaluate(_MEASURE_JS)
    if m["problems"]:
        raise Refusal("VOCAB_REFUSED", "the deck has content the PPTX emitter would "
                      "drop:\n  - " + "\n  - ".join(m["problems"][:8]))
    for i, s in enumerate(m["slides"], 1):
        if abs(s["w"] - 1280) > 1 or abs(s["h"] - 720) > 1:
            raise Refusal("GEOMETRY_UNDECLARED",
                          f"slide {i} measures {s['w']:.0f}x{s['h']:.0f} CSS px; a "
                          "deck slide is exactly 1280x720 (960x540pt)")
    return m["slides"]


def _write_text(slide, t: dict, font: str) -> None:
    from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
    wrap = t["full"] or t["lines"] > 1
    # a shrink-to-fit single-liner gets slack in the direction its alignment does not
    # anchor, so a substituted face overflows invisibly instead of re-wrapping
    width = t["w"] if wrap or t["align"] == "right" else t["w"] * LINE_SLACK
    tb = slide.shapes.add_textbox(_px(t["x"]), _px(t["y"]), _px(width), _px(t["h"]))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = {"right": PP_ALIGN.RIGHT, "center": PP_ALIGN.CENTER}.get(
        t["align"], PP_ALIGN.LEFT)
    if t.get("lh"):
        p.line_spacing = _pt(t["lh"])       # exact points, never a multiple
    run = p.add_run()
    run.text = shown(t)
    f = run.font
    f.name, f.size, f.bold = font, _pt(t["size"]), t["weight"] >= 600
    f.color.rgb = _rgb(t["colour"])
    if t.get("spacing"):
        run.font._rPr.set("spc", str(int(round(t["spacing"] * 0.75 * 100))))


def _write_table(slide, tbl: dict, font: str, slide_no: int, slide_h: float) -> None:
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    rows = tbl["rows"]
    if not rows:
        return
    total = sum(r["h"] for r in rows)
    # the projection gate: a table that fits Chrome and not PowerPoint refuses at
    # build time, not four days later when a human opens the file
    if tbl["y"] + total * PPTX_ROW_GROWTH > slide_h + 1:
        raise Refusal(
            "GATE_FAILED",
            f"slide {slide_no}: this table fits the browser and will not fit "
            f"PowerPoint: {len(rows)} rows measure {total:.0f}px in Chrome, about "
            f"{total * PPTX_ROW_GROWTH:.0f}px rendered, running past the slide "
            f"bottom. Take rows out, split the table, or take its type down -- do "
            "NOT raise the growth factor, which is measured")
    ncol = max(len(r["cells"]) for r in rows)
    shape = slide.shapes.add_table(len(rows), ncol, _px(tbl["x"]), _px(tbl["y"]),
                                   _px(tbl["w"]), _px(total))
    t = shape.table
    t.first_row = False
    t.horz_banding = False
    widths = [c["w"] for c in rows[0]["cells"]]
    widths += [0.0] * (ncol - len(widths))
    scale = tbl["w"] / max(sum(widths), 1.0)
    for i, w in enumerate(widths):
        t.columns[i].width = _px(max(w * scale, 8))
    for ri, row in enumerate(rows):
        t.rows[ri].height = _px(row["h"])
        for ci in range(ncol):
            cell = t.cell(ri, ci)
            cell.margin_left = _px(4)
            cell.margin_right = _px(4)
            cell.margin_top = cell.margin_bottom = _px(0)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if ci >= len(row["cells"]):
                cell.fill.background()
                continue
            spec = row["cells"][ci]
            fill = spec.get("fill", "")
            if fill and fill not in ("rgba(0, 0, 0, 0)", "transparent"):
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(fill)
            else:
                cell.fill.background()
            tf = cell.text_frame
            tf.word_wrap = True
            para = tf.paragraphs[0]
            para.alignment = {"right": PP_ALIGN.RIGHT,
                              "center": PP_ALIGN.CENTER}.get(spec["align"],
                                                             PP_ALIGN.LEFT)
            if spec["text"]:
                run = para.add_run()
                run.text = shown(spec)
                run.font.name = font
                run.font.size = _pt(spec["size"])
                run.font.bold = spec["weight"] >= 600
                run.font.color.rgb = _rgb(spec["colour"])
                if spec.get("spacing"):
                    run.font._rPr.set(
                        "spc", str(int(round(spec["spacing"] * 0.75 * 100))))


def build(slides_m: list[dict], shots: list[str], out_path: Path, font: str) -> int:
    from pptx import Presentation
    prs = Presentation()
    prs.slide_width, prs.slide_height = _px(1280), _px(720)
    blank = prs.slide_layouts[6]
    for i, sm in enumerate(slides_m, 1):
        slide = prs.slides.add_slide(blank)
        for f in sm["fills"]:
            sh = slide.shapes.add_shape(1, _px(f["x"]), _px(f["y"]),
                                        _px(f["w"]), _px(max(f["h"], 0.75)))
            sh.fill.solid()
            sh.fill.fore_color.rgb = _rgb(f["colour"])
            sh.line.fill.background()
            sh.shadow.inherit = False
        for img in sm["images"]:
            slide.shapes.add_picture(shots[img["idx"]], _px(img["x"]), _px(img["y"]),
                                     width=_px(img["w"]), height=_px(img["h"]))
        for tbl in sm["tables"]:
            _write_table(slide, tbl, font, i, sm["h"])
        for t in sm["texts"]:
            if t["text"]:
                _write_text(slide, t, font)
    prs.save(str(out_path))
    return len(slides_m)


def expected_runs(slides_m: list[dict]) -> list[str]:
    """Every string the browser painted, slide order irrelevant."""
    out = []
    for sm in slides_m:
        out += [shown(t) for t in sm["texts"] if t["text"]]
        for tbl in sm["tables"]:
            for row in tbl["rows"]:
                out += [shown(c) for c in row["cells"] if c["text"]]
    return out


def saved_runs(path: Path) -> list[list[str]]:
    """Text runs per slide, read off the SAVED bytes -- what ships is the file.
    Includes native table cells; a rasterised deck returns empty lists."""
    from pptx import Presentation
    out = []
    for slide in Presentation(str(path)).slides:
        runs = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    runs += [r.text for r in para.runs if r.text.strip()]
            if getattr(shape, "has_table", False) and shape.has_table:
                for cell in shape.table.iter_cells():
                    for para in cell.text_frame.paragraphs:
                        runs += [r.text for r in para.runs if r.text.strip()]
        out.append(runs)
    return out
